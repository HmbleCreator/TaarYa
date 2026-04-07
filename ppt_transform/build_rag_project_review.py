from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(r"C:\Users\amiku\Downloads\8th_Sem\TaarYa")
TEMPLATE = ROOT / "ppt_inspect" / "project-ppt-sample-converted.pptx"
SOURCE = ROOT / "RAG_Astronomy_Review.pptx"
WORKDIR = ROOT / "ppt_transform"
ASSET_DIR = WORKDIR / "assets"
OUTPUT = WORKDIR / "RAG_Astronomy_Review_Transformed_Project_Template.pptx"

FONT = "Times New Roman"
BLACK = RGBColor(0, 0, 0)


def delete_slide(prs: Presentation, index: int) -> None:
    slide_id = prs.slides._sldIdLst[index]
    rel_id = slide_id.rId
    prs.part.drop_rel(rel_id)
    del prs.slides._sldIdLst[index]


def set_run_font(run, size: int, bold: bool = False, italic: bool = False) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = BLACK


def disable_bullets(paragraph) -> None:
    p_pr = paragraph._p.get_or_add_pPr()
    for child in list(p_pr):
        if child.tag.endswith("}buNone") or child.tag.endswith("}buChar") or child.tag.endswith("}buAutoNum"):
            p_pr.remove(child)
    p_pr.insert(0, OxmlElement("a:buNone"))


def style_paragraph(paragraph, size: int, bold: bool = False, alignment=PP_ALIGN.LEFT, bullets: bool = False) -> None:
    paragraph.alignment = alignment
    paragraph.space_after = Pt(4)
    paragraph.space_before = Pt(0)
    paragraph.line_spacing = 1.1
    if not bullets:
        disable_bullets(paragraph)
    for run in paragraph.runs:
        set_run_font(run, size=size, bold=bold)


def set_title(slide, text: str, size: int = 32) -> None:
    title = slide.shapes.title
    title.text = text
    title_tf = title.text_frame
    title_tf.word_wrap = True
    title_tf.vertical_anchor = MSO_ANCHOR.TOP
    for paragraph in title_tf.paragraphs:
        style_paragraph(paragraph, size=size, bold=True, bullets=False)


def fill_text_placeholder(shape, lines, size: int = 20) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for idx, line in enumerate(lines):
        if idx == 0:
            paragraph = tf.paragraphs[0]
        else:
            paragraph = tf.add_paragraph()
        paragraph.text = line
        style_paragraph(paragraph, size=size, bullets=False)


def add_textbox(slide, left, top, width, height, lines, size=18, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for idx, line in enumerate(lines):
        if idx == 0:
            paragraph = tf.paragraphs[0]
        else:
            paragraph = tf.add_paragraph()
        paragraph.text = line
        style_paragraph(paragraph, size=size, bold=bold, alignment=align, bullets=False)
    return shape


def export_picture(shape, destination: Path) -> None:
    destination.write_bytes(shape.image.blob)


def extract_assets(source: Presentation) -> dict[str, Path]:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)

    diagrams = [shape for shape in source.slides[9].shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    diagram_names = [
        "class_diagram.png",
        "use_case_diagram.png",
        "dfd_level_0.png",
        "dfd_level_1.png",
        "dfd_level_2.png",
        "er_diagram.png",
    ]
    asset_map: dict[str, Path] = {}
    for name, shape in zip(diagram_names, diagrams):
        path = ASSET_DIR / name
        export_picture(shape, path)
        asset_map[name] = path

    gantt = next(shape for shape in source.slides[10].shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
    gantt_path = ASSET_DIR / "gantt_chart.png"
    export_picture(gantt, gantt_path)
    asset_map["gantt_chart.png"] = gantt_path

    screenshots = [shape for shape in source.slides[12].shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    screenshot_names = [
        "chat_interface.png",
        "query_results.png",
        "spatial_explorer.png",
        "system_status.png",
    ]
    for name, shape in zip(screenshot_names, screenshots):
        path = ASSET_DIR / name
        export_picture(shape, path)
        asset_map[name] = path

    return asset_map


def configure_cover_slide(slide) -> None:
    set_title(
        slide,
        "DEVELOPING AN INTELLIGENT\nRAG-DRIVEN ARCHITECTURE\nFOR ASTRONOMICAL STAR CATALOGS",
        size=30,
    )

    subtitle = slide.placeholders[1]
    subtitle.text = ""

    add_textbox(
        slide,
        Inches(4.2),
        Inches(3.55),
        Inches(4.8),
        Inches(0.45),
        ["Presented By"],
        size=22,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(1.55),
        Inches(4.4),
        Inches(3.8),
        Inches(1.2),
        ["Amit Kumar", "Roll No: 14618002722"],
        size=16,
        bold=True,
        align=PP_ALIGN.LEFT,
    )
    add_textbox(
        slide,
        Inches(8.2),
        Inches(4.35),
        Inches(3.1),
        Inches(1.45),
        ["Under Guidance of", "Dr. Madhumita Mahapatra", "Professor, CSE"],
        size=16,
        bold=False,
        align=PP_ALIGN.LEFT,
    )
    add_textbox(
        slide,
        Inches(4.1),
        Inches(6.55),
        Inches(5.1),
        Inches(0.3),
        ["Session 2025-2026 (Even SEM)"],
        size=14,
        align=PP_ALIGN.CENTER,
    )


def add_text_slide(prs: Presentation, title: str, lines, size: int = 20) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[11])
    set_title(slide, title)
    body = slide.placeholders[1]
    fill_text_placeholder(body, lines, size=size)


def configure_contents_slide(slide) -> None:
    lines = [
        "1. Introduction",
        "2. Problem Definition",
        "3. Feasibility Study",
        "4. Need and Significance",
        "5. Objectives",
        "6. Literature Survey",
        "7. Design Analysis / Diagrams",
        "8. Proposed Methodology in Brief",
        "9. Hardware and Software Requirements",
        "10. Running Project and Implementation",
        "11. Screen shots",
        "12. Scope of the Project",
        "13. Conclusion and Future Work",
        "14. References",
    ]
    set_title(slide, "Content")
    body = slide.placeholders[1]
    fill_text_placeholder(body, lines, size=18)


def add_diagram_slide(prs: Presentation, assets: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_title(slide, "Design Analysis / Diagrams")

    captions = [
        ("Class Diagram", assets["class_diagram.png"]),
        ("Use Case Diagram", assets["use_case_diagram.png"]),
        ("DFD Level 0", assets["dfd_level_0.png"]),
        ("DFD Level 1", assets["dfd_level_1.png"]),
        ("DFD Level 2", assets["dfd_level_2.png"]),
        ("E-R Diagram", assets["er_diagram.png"]),
    ]

    x_positions = [Inches(0.95), Inches(4.25), Inches(7.55)]
    y_positions = [Inches(1.35), Inches(4.05)]
    image_w = Inches(2.7)
    image_h = Inches(1.5)

    for idx, (caption, path) in enumerate(captions):
        col = idx % 3
        row = idx // 3
        left = x_positions[col]
        top = y_positions[row]
        slide.shapes.add_picture(str(path), left, top, width=image_w, height=image_h)
        add_textbox(
            slide,
            left,
            top + image_h + Inches(0.05),
            image_w,
            Inches(0.3),
            [caption],
            size=14,
            bold=True,
            align=PP_ALIGN.CENTER,
        )


def add_implementation_slide(prs: Presentation, assets: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_title(slide, "Running Project and Implementation")
    slide.shapes.add_picture(
        str(assets["gantt_chart.png"]),
        Inches(0.95),
        Inches(1.25),
        width=Inches(10.55),
        height=Inches(3.75),
    )
    add_textbox(
        slide,
        Inches(1.0),
        Inches(5.2),
        Inches(4.9),
        Inches(1.2),
        [
            "Current Development Status",
            "• Data Ingestion Engine — 90%",
            "• Vector Store Integration — 80%",
            "• ADQL/SQL Query Generator — 60%",
        ],
        size=15,
    )
    add_textbox(
        slide,
        Inches(6.0),
        Inches(5.2),
        Inches(4.9),
        Inches(1.2),
        [
            "• Knowledge Graph (Neo4j) — 70%",
            "• Agentic Reasoning Core — 70%",
            "• UI / Frontend — 60%",
            "• Minimum review milestone target: 70% complete",
        ],
        size=15,
    )


def add_screenshot_slide(prs: Presentation, assets: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_title(slide, "Screen shots")

    shots = [
        ("Chat Interface", assets["chat_interface.png"]),
        ("Spatial Explorer", assets["spatial_explorer.png"]),
        ("Query Results", assets["query_results.png"]),
        ("System Status", assets["system_status.png"]),
    ]
    positions = [
        (Inches(0.85), Inches(1.35)),
        (Inches(6.05), Inches(1.35)),
        (Inches(0.85), Inches(4.0)),
        (Inches(6.05), Inches(4.0)),
    ]

    for (caption, path), (left, top) in zip(shots, positions):
        slide.shapes.add_picture(str(path), left, top, width=Inches(4.55), height=Inches(2.2))
        add_textbox(
            slide,
            left,
            top + Inches(2.22),
            Inches(4.55),
            Inches(0.25),
            [caption],
            size=13,
            bold=True,
            align=PP_ALIGN.CENTER,
        )


def add_references_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    set_title(slide, "References")
    left = slide.placeholders[1]
    right = slide.placeholders[2]
    fill_text_placeholder(
        left,
        [
            "[1] Gao, Y., et al. Retrieval-Augmented Generation for Large Language Models: A Comprehensive Review. arXiv:2312.10997, 2024.",
            "[2] Lewis, P., et al. Retrieval-Augmented Generation for Knowledge-Intensive NLP. NeurIPS, 2020.",
            "[3] Edelman, G., et al. GraphRAG: Towards Graph-Based Retrieval-Augmented Generation. arXiv:2404.16130, 2024.",
            "[4] The Digital Architecture of the Heavens: A Comprehensive Analysis of Star Catalogs. arXiv:2501.00000, 2025.",
            "[5] A Comprehensive Report on Reasoning Language Models (RLMs). arXiv:2512.24601, 2025.",
        ],
        size=13,
    )
    fill_text_placeholder(
        right,
        [
            "[6] Vaswani, A., et al. Attention Is All You Need. NeurIPS, 2017.",
            "[7] Brown, T., et al. Language Models are Few-Shot Learners. NeurIPS, 2020.",
            "[8] Koposov, S., et al. Q3C, a Quad Tree Cube. Astronomy and Computing, 2010.",
            "[9] Gaia Collaboration. Gaia Data Release 3 Documentation.",
        ],
        size=13,
    )


def add_thank_you_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_title(slide, "THANK YOU", size=34)
    subtitle = slide.placeholders[1]
    subtitle.text = ""
    add_textbox(
        slide,
        Inches(3.0),
        Inches(4.2),
        Inches(7.0),
        Inches(1.2),
        [
            "Amit Kumar  |  Roll No: 14618002722",
            "Guide: Dr. Madhumita Mahapatra",
            "Delhi Technical Campus, Greater Noida",
            "Session 2025-2026",
        ],
        size=17,
        align=PP_ALIGN.CENTER,
    )


def build_deck() -> None:
    WORKDIR.mkdir(parents=True, exist_ok=True)

    template = Presentation(str(TEMPLATE))
    source = Presentation(str(SOURCE))
    assets = extract_assets(source)

    configure_cover_slide(template.slides[0])
    configure_contents_slide(template.slides[1])
    delete_slide(template, 2)
    add_text_slide(
        template,
        "Introduction",
        [
            "Petabyte-scale sky surveys such as Gaia and Vera C. Rubin have made astronomy a data-rich science with billions of stellar records.",
            "Researchers still work across siloed systems, switching manually between ADQL/SQL-based catalog search and keyword-based literature search.",
            "Retrieval-Augmented Generation (RAG) can unify these workflows by grounding large language models in astronomical catalogs and research papers.",
            "The proposed system acts as an intelligent co-pilot capable of natural-language querying, ADQL execution, and multi-step reasoning.",
            "Key data sources: Gaia DR3, SIMBAD, curated ArXiv astrophysics papers, and Astropy-powered analysis tools.",
        ],
        size=18,
    )
    add_text_slide(
        template,
        "Problem Definition",
        [
            "No unified interface exists for querying SIMBAD, Gaia, and ArXiv together in one research workflow.",
            "Students and early-career researchers face a technical barrier because ADQL/SQL tooling is rigid and difficult to learn.",
            "Standard semantic chunking misses the geometric and spatial relationships needed for cone searches and astronomical reasoning.",
            "Traditional RAG pipelines struggle with relational and multi-hop reasoning across observations, papers, and catalog entities.",
            "Current systems lack agentic capabilities for planning research steps, selecting tools, and verifying answers against source data.",
        ],
        size=18,
    )
    add_text_slide(
        template,
        "Feasibility Study",
        [
            "Technical Feasibility",
            "LangChain and LlamaIndex support agentic orchestration.",
            "Astropy and astroquery support Gaia and SIMBAD integration.",
            "Neo4j and PostgreSQL with Q3C provide graph and spatial indexing at scale.",
            "Economic / Operational Feasibility",
            "Open-source models and public datasets keep the solution cost-effective and practical for academic deployment.",
            "A web-based modular architecture supports incremental rollout from core search to full agentic features.",
            "Schedule Feasibility",
            "Month 1: literature review and connectors; Month 2: ingestion and indexing; Month 3: graph and agentic core; Month 4: integration and evaluation.",
        ],
        size=17,
    )
    add_text_slide(
        template,
        "Need and Significance",
        [
            "Modern astroinformatics produces data volumes too large for manual exploration, making intelligent retrieval essential.",
            "For researchers: the system creates a unified data-to-discovery interface that links catalogs and literature in one workflow.",
            "For education: students can work in plain English instead of depending on advanced ADQL or scripting expertise.",
            "For the field: the project combines SQL precision, vector semantics, and graph logic into a scientific AI methodology.",
            "As a future-facing contribution, the architecture aligns with Reasoning Language Models that actively navigate data environments.",
        ],
        size=18,
    )
    add_text_slide(
        template,
        "Objectives",
        [
            "1. Build a multi-modal ingestion engine for FITS/CSV catalogs, ArXiv PDFs, and observational metadata.",
            "2. Create a hybrid indexing system using dense vectors, BM25 keyword retrieval, and Q3C spatial indexing.",
            "3. Construct a knowledge graph linking stars, papers, catalogs, and cross-matched entities.",
            "4. Develop an agentic reasoning core for query decomposition, tool selection, and execution.",
            "5. Enable code execution with Astropy and Pandas for on-the-fly analysis and plotting.",
            "6. Add a verification and synthesis layer that checks answers against catalog values and produces cited responses.",
            "7. Evaluate the system with astronomy-adapted RAG benchmarks and retrieval metrics.",
        ],
        size=17,
    )
    add_text_slide(
        template,
        "Literature Survey",
        [
            "Lewis et al. established the foundational RAG framework for grounding LLMs in external knowledge bases.",
            "Gao et al. organized RAG systems into a progression from naive to advanced, modular, and agentic architectures.",
            "Edelman et al. introduced GraphRAG to capture relational structure that vector-only retrieval often misses.",
            "Prior star-catalog work shows the shift from positional lists to multidimensional databases with Q3C-style spatial indexing.",
            "Recent Reasoning Language Model work highlights iterative tool use and deliberate search strategies for complex tasks.",
            "Proposed improvements: agentic tool use, hybrid indexing, graph-enhanced reasoning, and Astropy-driven interaction with astronomy data.",
        ],
        size=17,
    )
    add_diagram_slide(template, assets)
    add_text_slide(
        template,
        "Proposed Methodology in Brief",
        [
            "M1: Multi-Modal Ingestion Engine for FITS/CSV files, ArXiv PDFs, and entity extraction.",
            "M2: Hybrid Indexing System combining dense vectors, BM25 retrieval, and Q3C spatial indexing.",
            "M3: Knowledge Graph Constructor with nodes such as stars, papers, and catalogs, plus cross-reference edges.",
            "M4: Agentic Reasoning Core for query planning, tool routing, and execution.",
            "M5: Verification and Synthesis Layer for hallucination checking and citation generation.",
            "Non-functional goals: ADQL/SQL latency below 3 seconds, precise cone-search accuracy, and a student-friendly natural-language UI.",
        ],
        size=17,
    )
    add_text_slide(
        template,
        "Hardware and Software Requirements",
        [
            "Hardware Requirements",
            "• Server RAM: 16 GB minimum for vector database operations.",
            "• GPU: NVIDIA T4 optional for local LLM hosting and experimentation.",
            "• Storage: 500 GB SSD for Gaia catalog subsets and intermediate artifacts.",
            "Software Requirements",
            "• Backend: Python 3.9+ with FastAPI.",
            "• LLM Layer: LangChain / LlamaIndex.",
            "• Databases: Neo4j, PostgreSQL + Q3C, and Qdrant.",
            "• Libraries: Astropy, Pandas, Matplotlib, and astroquery.",
        ],
        size=17,
    )
    add_implementation_slide(template, assets)
    add_screenshot_slide(template, assets)
    add_text_slide(
        template,
        "Scope of the Project",
        [
            "Domain scope includes Gaia DR3 catalogs, SIMBAD records, and a curated set of astrophysics papers from ArXiv.",
            "Functional scope covers natural-language querying, catalog lookup, literature synthesis, hybrid retrieval, graph reasoning, and result export.",
            "Analytical scope includes spatial search, cross-referencing of catalog entities, and code-assisted statistical analysis through Astropy/Pandas.",
            "Deployment scope is a web-based academic research assistant designed for researchers and students in astroinformatics workflows.",
        ],
        size=18,
    )
    add_text_slide(
        template,
        "Conclusion and Future Work",
        [
            "The proposed RAG-driven architecture addresses a real astronomy workflow gap by unifying data retrieval, literature search, and reasoning.",
            "A hybrid approach combining vector retrieval, spatial indexing, knowledge graphs, and agentic tool use is well-suited to star-catalog research.",
            "Prototype progress already demonstrates strong momentum in ingestion, retrieval, graph construction, and interface development.",
            "Future work includes completing the ADQL generation pipeline, expanding benchmark evaluation, broadening catalog coverage, and refining the end-user interface.",
        ],
        size=18,
    )
    add_references_slide(template)
    add_thank_you_slide(template)

    template.save(str(OUTPUT))


if __name__ == "__main__":
    build_deck()
