"""
Extract .ppt content to HTML using python-pptx.
Note: python-pptx natively supports .pptx; .ppt (binary) files need conversion first.
We'll try to open it directly (some .ppt files masquerade as .pptx internally).
"""
import sys
import os

ppt_path = r"c:\Users\amiku\Downloads\8th_Sem\TaarYa\project ppt sample.ppt"
out_path = r"c:\Users\amiku\Downloads\8th_Sem\TaarYa\tmp_ppt_content.html"

try:
    from pptx import Presentation
    from pptx.util import Pt
    
    prs = Presentation(ppt_path)
    
    html_parts = ["""<!DOCTYPE html>
<html>
<head><meta charset='utf-8'><title>PPT Content</title>
<style>
body { font-family: Arial, sans-serif; margin: 20px; }
.slide { border: 2px solid #333; margin: 30px 0; padding: 20px; background: #f9f9f9; }
.slide-num { font-size: 12px; color: #888; margin-bottom: 10px; }
.title { font-size: 24px; font-weight: bold; margin-bottom: 10px; color: #1a1a2e; }
.content { font-size: 14px; margin: 5px 0; }
.bullet { margin-left: 20px; }
.shape-name { font-size: 10px; color: #aaa; margin-top: 5px; font-style: italic; }
</style>
</head><body>
"""]
    
    for slide_num, slide in enumerate(prs.slides, 1):
        layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
        html_parts.append(f'<div class="slide">')
        html_parts.append(f'<div class="slide-num">SLIDE {slide_num} | Layout: {layout_name}</div>')
        
        for shape in slide.shapes:
            shape_type = str(shape.shape_type)
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                shape_name = shape.name if shape.name else "unnamed"
                lines = text.split("\n")
                
                # Check if it's likely a title
                if "title" in shape_name.lower() or slide_num == 1 and shape == slide.shapes[0]:
                    html_parts.append(f'<div class="title">{lines[0]}</div>')
                    for line in lines[1:]:
                        if line.strip():
                            html_parts.append(f'<div class="content">{line}</div>')
                else:
                    for line in lines:
                        if line.strip():
                            # detect bullets
                            if line.startswith("•") or line.startswith("-") or line.startswith("*"):
                                html_parts.append(f'<div class="bullet">{line}</div>')
                            else:
                                html_parts.append(f'<div class="content">{line}</div>')
                
                html_parts.append(f'<div class="shape-name">[Shape: {shape_name} | Type: {shape_type}]</div>')
        
        # Check for notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                html_parts.append(f'<div class="content" style="color:#555;border-top:1px solid #ddd;margin-top:10px;"><b>Notes:</b> {notes}</div>')
        
        html_parts.append('</div>')
    
    html_parts.append("</body></html>")
    
    with open(out_path, "w", encoding="utf-8") as f:
        f.write("\n".join(html_parts))
    
    print(f"SUCCESS: Extracted {len(prs.slides)} slides to {out_path}")

except Exception as e:
    print(f"ERROR: {e}")
    # Try installing python-pptx if missing
    if "No module named" in str(e):
        print("Try: pip install python-pptx")
    sys.exit(1)
